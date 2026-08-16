import pdfplumber
from pptx import Presentation
from pathlib import Path
from typing import Dict, List
import re

def extract_pdf(path: str) -> Dict:
    pages: List[Dict] = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            raw_text = page.extract_text() or ""
            tables = page.extract_tables() or []
            
            pages.append({
                "page_number": i + 1,
                "raw_text": raw_text,
                "tables": tables
            })
    full_text = "\n".join(p["raw_text"] for p in pages)
    return {
            "source_type": "pdf",
            "page_count": len(pages),
            "pages": pages,
            "raw_text": full_text
            }

def extract_from_pptx(file_path: str) -> Dict:    
    prs = Presentation(file_path)
    slides: List[Dict] = []

    for i, slide in enumerate(prs.slides):
        slide_lines: List[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line:
                        slide_lines.append(line)

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text:
                            slide_lines.append(cell.text)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        slides.append({
            "slide_number": i + 1,
            "raw_text": "\n".join(slide_lines),
            "notes": notes
        })

    full_text = "\n".join(s["raw_text"] for s in slides)

    return {
        "source_type": "pptx",
        "slide_count": len(slides),
        "slides": slides,
        "raw_text": full_text
    }


def extract_text(file_path: str) -> Dict:

    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return extract_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def semantic_chunk(text: str, source_type: str = "pdf") -> List[Dict]:
    if source_type == "pptx":
        raw_chunks = text.split("\n\n")
    else:
        # split on heading-like pattern or double line breaks
        raw_chunks = re.split(r"\n(?=[A-Z][^\n]{0,80}\n)|\n{2,}", text)
    
    chunks = []
    for i, chunk in enumerate(raw_chunks):
        chunk = chunk.strip()
        if len(chunk) < 20:  # skip noise/empty fragments
            continue
        chunks.append({
            "chunk_id": i,
            "text": chunk,
            "char_len": len(chunk)
        })
    
    return chunks
